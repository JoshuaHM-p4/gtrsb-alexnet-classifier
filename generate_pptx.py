import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_base_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Solid background (light off-white)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 57, 107) # Dark Blue Primary
    
    # Thin decorative line below the title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(46, 125, 50) # Forest Green Secondary
    line.line.fill.background() # No border
    
    return slide

def add_bullet(tf, text, level=0, bold=False, size=16, color=RGBColor(51, 51, 51)):
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    
    # Apply clean bullet design by prepending bullet character
    p.text = f"•  {text}"
    p.level = level
    p.font.name = "Segoe UI"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(12) # Nice spacing between bullets
    p.space_before = Pt(4)

def set_cell_text(cell, text, bold=False, size=11, color=RGBColor(40, 40, 40), align=PP_ALIGN.LEFT, bg_color=None):
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    cell.text_frame.text = "" # clear default
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_full_image_slide(prs, title_text, img_path):
    slide = create_base_slide(prs, title_text)
    if not os.path.exists(img_path):
        print(f"Warning: Image {img_path} not found.")
        return slide
    
    # Calculate box sizes to center the image on the screen
    # Screen height is 7.5, title takes ~1.5, remaining height is 6.0
    # Screen width is 13.33
    img_width = Inches(8.5)
    img_height = Inches(5.2)
    left = Inches((13.33 - 8.5) / 2)
    top = Inches(1.7)
    
    # Add border shape
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left - Inches(0.04), top - Inches(0.04), img_width + Inches(0.08), img_height + Inches(0.08)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(220, 220, 220)
    border.line.fill.background()
    
    # Add picture
    slide.shapes.add_picture(img_path, left, top, img_width, img_height)
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide_1 = prs.slides.add_slide(blank_layout)
    bg = slide_1.background
    bg_fill = bg.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # Center PUP logo at top
    logo_path = "docu/icons/pup_logo.png"
    if os.path.exists(logo_path):
        slide_1.shapes.add_picture(logo_path, Inches(5.66), Inches(0.6), Inches(2.0), Inches(2.0))
        
    title_box = slide_1.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.33), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Image Classification Using CNN Architectures"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(21, 57, 107)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "AlexNet Classification on German Traffic Sign Recognition Benchmark (GTSRB)"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(20)
    p2.font.bold = False
    p2.font.color.rgb = RGBColor(46, 125, 50)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    
    info_box = slide_1.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.33), Inches(2.0))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    p3 = tf_info.paragraphs[0]
    p3.text = "POLYTECHNIC UNIVERSITY OF THE PHILIPPINES | College of Engineering"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(100, 100, 100)
    p3.alignment = PP_ALIGN.CENTER
    
    p4 = tf_info.add_paragraph()
    p4.text = "Presented by: Carlos Jerico S. Dela Torre, Joshua H. Mistal, Aidan R. Tiu"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(60, 60, 60)
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(10)
    
    p5 = tf_info.add_paragraph()
    p5.text = "Prof. Mon Arjay F. Malbog | CMPE 362 Pattern Recognition | June 2026"
    p5.font.name = "Segoe UI"
    p5.font.size = Pt(11)
    p5.font.color.rgb = RGBColor(120, 120, 120)
    p5.alignment = PP_ALIGN.CENTER
    p5.space_before = Pt(5)

    # ==========================================
    # SLIDE 2: Project Overview & Abstract
    # ==========================================
    slide_2 = create_base_slide(prs, "Project Overview & Executive Summary")
    
    txt_box = slide_2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(5.0))
    tf = txt_box.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "TSR Evaluation: Evaluates AlexNet on safety-critical traffic sign classification.", size=17, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf, "Transfer Learning Strategy: Freezes ImageNet layers; retrains modified final classification layers.")
    add_bullet(tf, "Overcoming Noisy Conditions: Successfully handles shadows, motion blur, and low lighting.")
    add_bullet(tf, "Imbalance Resilience: Mitigates massive class representation differences across the 43 target classes.")
    
    callout = slide_2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.08), Inches(4.8))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(240, 244, 248)
    callout.line.color.rgb = RGBColor(21, 57, 107)
    callout.line.width = Pt(1.5)
    
    tf_c = callout.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Inches(0.4)
    tf_c.margin_right = Inches(0.4)
    tf_c.margin_top = Inches(0.4)
    
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "Key Metrics on Test Set"
    p_c1.font.name = "Segoe UI"
    p_c1.font.size = Pt(22)
    p_c1.font.bold = True
    p_c1.font.color.rgb = RGBColor(21, 57, 107)
    p_c1.space_after = Pt(20)
    
    p_c2 = tf_c.add_paragraph()
    p_c2.text = "Test Accuracy:"
    p_c2.font.name = "Segoe UI"
    p_c2.font.size = Pt(14)
    p_c2.font.color.rgb = RGBColor(80, 80, 80)
    
    p_c3 = tf_c.add_paragraph()
    p_c3.text = "98.91%"
    p_c3.font.name = "Segoe UI"
    p_c3.font.size = Pt(40)
    p_c3.font.bold = True
    p_c3.font.color.rgb = RGBColor(46, 125, 50)
    p_c3.space_after = Pt(15)
    
    p_c4 = tf_c.add_paragraph()
    p_c4.text = "Precision (Macro): 0.9924\nRecall (Macro): 0.9899\nF1-Score (Macro): 0.9910\nROC-AUC (Macro): 1.0000"
    p_c4.font.name = "Segoe UI"
    p_c4.font.size = Pt(15)
    p_c4.font.bold = True
    p_c4.font.color.rgb = RGBColor(51, 51, 51)
    
    # ==========================================
    # SLIDE 3: Chapter 1: Introduction - Background & Problem
    # ==========================================
    slide_3 = create_base_slide(prs, "Chapter 1: Background & Problem Statement")
    
    box_left = slide_3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(5.0))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    add_bullet(tf_l, "Background Context", bold=True, size=20, color=RGBColor(21, 57, 107))
    add_bullet(tf_l, "TSR is crucial for Advanced Driver Assistance Systems (ADAS) and autonomous safety.")
    add_bullet(tf_l, "GTSRB contains over 50,000 traffic sign images across 43 classes.")
    add_bullet(tf_l, "Convolutional Networks replace manual color/shape extraction with automatic features.")
    
    box_right = slide_3.shapes.add_textbox(Inches(6.98), Inches(1.8), Inches(5.6), Inches(5.0))
    tf_r = box_right.text_frame
    tf_r.word_wrap = True
    add_bullet(tf_r, "Problem Definition", bold=True, size=20, color=RGBColor(198, 40, 40))
    add_bullet(tf_r, "Images suffer from severe blur, low lighting, and perspective distortion.")
    add_bullet(tf_r, "Highly imbalanced classes bias general classification networks.")
    add_bullet(tf_r, "Deploying deep models on vehicle edge computers requires balancing speed and accuracy.")

    # ==========================================
    # SLIDE 4: Chapter 1: Objectives
    # ==========================================
    slide_4 = create_base_slide(prs, "Chapter 1: Project Objectives")
    
    card_width = Inches(3.6)
    card_height = Inches(4.5)
    card_y = Inches(2.0)
    spacing = Inches(0.5)
    start_x = Inches(0.75)
    
    objectives = [
        ("1. Data Preprocessing", "Standardize dimension and lighting variables across inputs.", [
            "Resize all raw images to standard 256x256 grids.",
            "Center-crop to 224x224x3 to fit standard models.",
            "Apply ImageNet channel normalization."
        ]),
        ("2. Architecture Adaptation", "Modify convolutional networks for transfer learning.", [
            "Load pre-trained PyTorch AlexNet weights.",
            "Freeze feature extraction convolutional layers.",
            "Modify final classifier layer for 43 classes."
        ]),
        ("3. Model Evaluation", "Train adapted layers and compute validation metrics.", [
            "Minimize Cross Entropy Loss using Adam.",
            "Integrate early stopping (patience=3) to prevent overfitting.",
            "Evaluate test accuracy, precision, recall, and ROC-AUC."
        ])
    ]
    
    for idx, (title, subtitle, bullets) in enumerate(objectives):
        card_x = start_x + idx * (card_width + spacing)
        card = slide_4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 247, 250)
        card.line.color.rgb = RGBColor(21, 57, 107)
        card.line.width = Pt(1.5)
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = tf_card.margin_right = Inches(0.2)
        tf_card.margin_top = Inches(0.3)
        
        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Segoe UI"
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = RGBColor(21, 57, 107)
        p_t.space_after = Pt(5)
        
        p_s = tf_card.add_paragraph()
        p_s.text = subtitle
        p_s.font.name = "Segoe UI"
        p_s.font.size = Pt(12)
        p_s.font.italic = True
        p_s.font.color.rgb = RGBColor(100, 100, 100)
        p_s.space_after = Pt(15)
        
        for bullet in bullets:
            add_bullet(tf_card, bullet, level=0, size=11, color=RGBColor(60, 60, 60))

    # ==========================================
    # SLIDE 5: Chapter 2: Literature Review
    # ==========================================
    slide_5 = create_base_slide(prs, "Chapter 2: Literature Review")
    
    tf_l5 = slide_5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l5.word_wrap = True
    
    add_bullet(tf_l5, "Early Methods & Issues: Thresholding and shape descriptors failed under dynamic shadows (Stallkamp et al., 2012).", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l5, "The CNN Era: Sermanet and LeCun (2011) proved multi-scale CNNs outperform human visual accuracy on GTSRB.")
    add_bullet(tf_l5, "AlexNet Customization: Ma et al. (2018) optimized AlexNet for traffic signs by shrinking kernels to prevent overfitting.")
    add_bullet(tf_l5, "Performance-Speed Balance: Arcos-Garcia et al. (2025) confirmed AlexNet provides highly efficient speed-accuracy tradeoffs for edge vehicle systems.")

    # ==========================================
    # SLIDE 6: Chapter 2: CNN Architecture Overview
    # ==========================================
    slide_6 = create_base_slide(prs, "Chapter 2: CNN Architecture Overview")
    
    tf_l6 = slide_6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l6.word_wrap = True
    
    add_bullet(tf_l6, "Preservation of Spatial Data: Retains grid relationship structure through local receptive fields.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l6, "Weight Sharing Advantage: Shared kernel arrays reduce total parameters, accelerating training.")
    add_bullet(tf_l6, "Hierarchical Layer Layout: Features scale from simple lines/edges to semantic sign classes.")

    # ==========================================
    # SLIDE 7: Table: CNN Basic Layers
    # ==========================================
    slide_7 = create_base_slide(prs, "CNN Structural Layer Summary")
    
    # Programmatic table centered on slide
    rows, cols = 5, 2
    table_shape = slide_7.shapes.add_table(rows, cols, Inches(2.0), Inches(1.8), Inches(9.33), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(6.83)
    
    set_cell_text(table.cell(0, 0), "Layer Type", bold=True, size=15, color=RGBColor(255, 255, 255), bg_color=RGBColor(21, 57, 107))
    set_cell_text(table.cell(0, 1), "Primary Structural Function", bold=True, size=15, color=RGBColor(255, 255, 255), bg_color=RGBColor(21, 57, 107))
    
    layer_data = [
        ("Convolutional", "Convolves sliding filters with input to extract spatial feature maps."),
        ("Activation (ReLU)", "Introduces non-linearity, allowing complex classification boundary learning."),
        ("Pooling (Max)", "Reduces spatial dimensions, parameter count, and compute overhead."),
        ("Fully Connected", "Maps high-level abstracted feature lists to target classification labels.")
    ]
    
    for r_idx, (layer_type, desc) in enumerate(layer_data, start=1):
        bg_col = RGBColor(240, 242, 245) if r_idx % 2 == 0 else RGBColor(255, 255, 255)
        set_cell_text(table.cell(r_idx, 0), layer_type, bold=True, size=13, color=RGBColor(21, 57, 107), bg_color=bg_col)
        set_cell_text(table.cell(r_idx, 1), desc, bold=False, size=12, color=RGBColor(51, 51, 51), bg_color=bg_col)

    # ==========================================
    # SLIDE 8: Chapter 2: AlexNet Architecture Overview
    # ==========================================
    slide_8 = create_base_slide(prs, "Chapter 2: AlexNet Architecture Details")
    
    tf_l8 = slide_8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l8.word_wrap = True
    
    add_bullet(tf_l8, "Network Layout: Consists of 5 Convolutional layers and 3 Fully Connected layers.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l8, "PyTorch Input Grid: Standardized to 224x224x3 using a padding of 2 in the first layer.")
    add_bullet(tf_l8, "Dropout Regularization: Uses Dropout layers (p = 0.5) to prevent unit co-adaptation and overfitting.")
    add_bullet(tf_l8, "Softmax Integration: Logits are parsed to output 43 categorical probabilities.")

    # ==========================================
    # SLIDE 9: Figure: AlexNet Block Diagram (DEDICATED SLIDE)
    # ==========================================
    add_full_image_slide(prs, "Figure 1: AlexNet Architectural Layout Diagram", "docu/figures/alexnet_diagram.jpg")

    # ==========================================
    # SLIDE 10: Chapter 3: Methodology Overview
    # ==========================================
    slide_10 = create_base_slide(prs, "Chapter 3: Methodology Overview")
    
    tf_l10 = slide_10.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l10.word_wrap = True
    
    add_bullet(tf_l10, "Development Frameworks: Implemented in Python 3 with CUDA GPU hardware acceleration.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l10, "Deep Learning Pipeline: Built using PyTorch and Torchvision for network graphing.")
    add_bullet(tf_l10, "Metrics & Splits: Scikit-learn executed stratified splits and computed macro classification scores.")
    add_bullet(tf_l10, "Data Visualization: Matplotlib and Seaborn mapped training loss/accuracy history.")

    # ==========================================
    # SLIDE 11: Figure: GTSRB Sample Signs (DEDICATED SLIDE)
    # ==========================================
    add_full_image_slide(prs, "Figure 2: Sample Signs Across all 43 GTSRB Classes", "docu/figures/sample_all_classes.png")

    # ==========================================
    # SLIDE 12: Figure: Dataset Sample Distribution (DEDICATED SLIDE)
    # ==========================================
    add_full_image_slide(prs, "Figure 3: Class-wise Sample Distribution Heatmap", "docu/figures/dataset_sample_distributions.png")

    # ==========================================
    # SLIDE 13: Chapter 3: Data Loading & Split
    # ==========================================
    slide_13 = create_base_slide(prs, "Chapter 3: Data Loading & Split")
    
    tf_l13 = slide_13.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l13.word_wrap = True
    
    add_bullet(tf_l13, "Stratified Split Strategy: Unified all train/test paths, then split into 70% Train, 15% Val, 15% Test.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l13, "Stratification: Preserves minority class proportions across all three subsets to prevent classifier bias.")
    add_bullet(tf_l13, "DataLoader Acceleration: Batches of 32 utilizing multi-process worker pools (num_workers=2) for high GPU throughput.")

    # ==========================================
    # SLIDE 14: Code: Custom GTSRBDataset
    # ==========================================
    slide_14 = create_base_slide(prs, "Methodology: PyTorch Dataset Class")
    
    code_card = slide_14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.8), Inches(9.33), Inches(4.8))
    code_card.fill.solid()
    code_card.fill.fore_color.rgb = RGBColor(240, 240, 240)
    code_card.line.color.rgb = RGBColor(200, 200, 200)
    
    tf_code = code_card.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = Inches(0.3)
    tf_code.margin_top = Inches(0.3)
    
    code_p = tf_code.paragraphs[0]
    code_p.text = "Custom Dataset Implementation (GTSRBDataset)"
    code_p.font.name = "Segoe UI"
    code_p.font.size = Pt(14)
    code_p.font.bold = True
    code_p.font.color.rgb = RGBColor(21, 57, 107)
    code_p.space_after = Pt(10)
    
    code_text = (
        "class GTSRBDataset(Dataset):\n"
        "    def __init__(self, dataframe, root_dir, transform=None):\n"
        "        self.dataframe = dataframe.reset_index(drop=True)\n"
        "        self.root_dir = root_dir\n"
        "        self.transform = transform\n"
        "        self.image_paths = self.dataframe['Path'].values\n"
        "        self.labels = self.dataframe['ClassId'].values\n\n"
        "    def __getitem__(self, idx):\n"
        "        img_path = os.path.join(self.root_dir, self.image_paths[idx])\n"
        "        image = Image.open(img_path).convert('RGB')\n"
        "        label = self.labels[idx]\n"
        "        if self.transform: image = self.transform(image)\n"
        "        return image, label"
    )
    
    p_code = tf_code.add_paragraph()
    p_code.text = code_text
    p_code.font.name = "Consolas"
    p_code.font.size = Pt(11)
    p_code.font.color.rgb = RGBColor(40, 40, 40)

    # ==========================================
    # SLIDE 15: Chapter 3: Preprocessing & Data Augmentation
    # ==========================================
    slide_15 = create_base_slide(prs, "Chapter 3: Preprocessing & Augmentation")
    
    box_left = slide_15.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(5.0))
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    add_bullet(tf_l, "Training Augmentation Pipeline", bold=True, size=18, color=RGBColor(21, 57, 107))
    add_bullet(tf_l, "Resizes inputs to a standardized 256x256 pixel grid.")
    add_bullet(tf_l, "Random rotation of up to 15° simulates camera tilt and perspectives.")
    add_bullet(tf_l, "Color Jitter (0.2 factor) alters brightness, contrast, and saturation.")
    add_bullet(tf_l, "Center crop to 224x224 and normalize via ImageNet values.")
    
    box_right = slide_15.shapes.add_textbox(Inches(6.98), Inches(1.8), Inches(5.6), Inches(5.0))
    tf_r = box_right.text_frame
    tf_r.word_wrap = True
    add_bullet(tf_r, "Validation & Test Pipeline (Static)", bold=True, size=18, color=RGBColor(46, 125, 50))
    add_bullet(tf_r, "Resizes inputs to a standardized 256x256 pixel grid.")
    add_bullet(tf_r, "Omits random rotations and color jittering to prevent evaluation variance.")
    add_bullet(tf_r, "Applies identical center cropping to 224x224 and ImageNet normalizations.")

    # ==========================================
    # SLIDE 16: Chapter 3: Transfer Learning Setup
    # ==========================================
    slide_16 = create_base_slide(prs, "Chapter 3: Transfer Learning Setup")
    
    tf_l16 = slide_16.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l16.word_wrap = True
    
    add_bullet(tf_l16, "Frozen Feature Extractors: Disables gradients (`requires_grad=False`) on convolutional layers.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l16, "Pre-trained ImageNet Weights: Retains low-level descriptors (edges, geometry) without training overhead.")
    add_bullet(tf_l16, "Classification Head Swap: Replaces index 6 linear layer to output 43 categorical boundaries.")

    # ==========================================
    # SLIDE 17: Code: Transfer Learning Setup
    # ==========================================
    slide_17 = create_base_slide(prs, "Methodology: Model Initialization")
    
    code_card2 = slide_17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.8), Inches(9.33), Inches(4.8))
    code_card2.fill.solid()
    code_card2.fill.fore_color.rgb = RGBColor(240, 240, 240)
    code_card2.line.color.rgb = RGBColor(200, 200, 200)
    
    tf_code2 = code_card2.text_frame
    tf_code2.word_wrap = True
    tf_code2.margin_left = Inches(0.3)
    tf_code2.margin_top = Inches(0.3)
    
    code_p2 = tf_code2.paragraphs[0]
    code_p2.text = "Model Customization and Freezing (PyTorch)"
    code_p2.font.name = "Segoe UI"
    code_p2.font.size = Pt(14)
    code_p2.font.bold = True
    code_p2.font.color.rgb = RGBColor(21, 57, 107)
    code_p2.space_after = Pt(10)
    
    code_text2 = (
        "# Load pre-trained AlexNet\n"
        "model = models.alexnet(weights=AlexNet_Weights.DEFAULT)\n\n"
        "# Freeze convolutional feature layers\n"
        "for param in model.features.parameters():\n"
        "    param.requires_grad = False\n\n"
        "# Reconstruct classification layer for 43 classes\n"
        "num_ftrs = model.classifier[6].in_features\n"
        "model.classifier[6] = nn.Linear(num_ftrs, 43)\n\n"
        "# Push parameters to CUDA GPU\n"
        "model = model.to(device)"
    )
    
    p_code2 = tf_code2.add_paragraph()
    p_code2.text = code_text2
    p_code2.font.name = "Consolas"
    p_code2.font.size = Pt(11)
    p_code2.font.color.rgb = RGBColor(40, 40, 40)

    # ==========================================
    # SLIDE 18: Chapter 3: Training Process Configuration
    # ==========================================
    slide_18 = create_base_slide(prs, "Chapter 3: Training & Optimization Settings")
    
    tf_l18 = slide_18.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l18.word_wrap = True
    
    add_bullet(tf_l18, "Loss Criteria: Minimizes Cross Entropy Loss, mapping multi-class categorical distributions.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l18, "Optimization settings: Adam optimizer restricted to `model.classifier.parameters()` (learning rate = 0.0001).")
    add_bullet(tf_l18, "Early Stopping Regulation: Patience of 3 epochs stops training early if validation loss fails to decrease.")
    add_bullet(tf_l18, "Max Epochs: Set to 15 epochs; restores the best weights after termination.")

    # ==========================================
    # SLIDE 19: Chapter 4: Results & Performance Metrics
    # ==========================================
    slide_19 = create_base_slide(prs, "Chapter 4: Results - Classification Summary")
    
    tf_l19 = slide_19.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l19.word_wrap = True
    
    add_bullet(tf_l19, "Testing set evaluation: The model was evaluated on 15% of the pooled dataset.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l19, "High Sensitivity and Specificity: Macro Precision of 0.9924 and Recall of 0.9899 confirm strong performance.")
    add_bullet(tf_l19, "F1-Score Balance: Macro F1-score of 0.9910 indicates that minority traffic signs are predicted reliably.")

    # ==========================================
    # SLIDE 20: Table: Classification Performance
    # ==========================================
    slide_20 = create_base_slide(prs, "Overall Performance Score Sheet")
    
    rows, cols = 6, 2
    table_shape2 = slide_20.shapes.add_table(rows, cols, Inches(2.0), Inches(1.8), Inches(9.33), Inches(4.5))
    table2 = table_shape2.table
    table2.columns[0].width = Inches(5.0)
    table2.columns[1].width = Inches(4.33)
    
    set_cell_text(table2.cell(0, 0), "Evaluation Metric", bold=True, size=15, color=RGBColor(255, 255, 255), bg_color=RGBColor(21, 57, 107))
    set_cell_text(table2.cell(0, 1), "Test Set Score", bold=True, size=15, color=RGBColor(255, 255, 255), bg_color=RGBColor(21, 57, 107))
    
    metrics_data = [
        ("Test Accuracy", "98.91%"),
        ("Precision (Macro Avg)", "0.9924"),
        ("Recall (Macro Avg)", "0.9899"),
        ("F1-Score (Macro Avg)", "0.9910"),
        ("ROC-AUC (Macro Avg)", "1.0000")
    ]
    
    for r_idx, (metric, value) in enumerate(metrics_data, start=1):
        bg_col = RGBColor(240, 242, 245) if r_idx % 2 == 0 else RGBColor(255, 255, 255)
        set_cell_text(table2.cell(r_idx, 0), metric, bold=True, size=13, color=RGBColor(21, 57, 107), bg_color=bg_col)
        val_color = RGBColor(46, 125, 50) if "%" in value or "1.0000" in value else RGBColor(51, 51, 51)
        set_cell_text(table2.cell(r_idx, 1), value, bold=True, size=13, color=val_color, bg_color=bg_col, align=PP_ALIGN.CENTER)

    # ==========================================
    # SLIDE 21: Figure: Loss & Accuracy Curves (DEDICATED SLIDE)
    # ==========================================
    add_full_image_slide(prs, "Figure 4: Training and Validation Loss/Accuracy History", "docu/figures/loss_acc_per_epoch.png")

    # ==========================================
    # SLIDE 22: Figure: Confusion Matrix (DEDICATED SLIDE)
    # ==========================================
    add_full_image_slide(prs, "Figure 5: Confusion Matrix for 43 Classes", "docu/figures/accuracy_graph.png")

    # ==========================================
    # SLIDE 23: Figure: ROC-AUC Analysis (DEDICATED SLIDE)
    # ==========================================
    add_full_image_slide(prs, "Figure 6: Macro-Average Receiver Operating Curve (ROC)", "docu/figures/roc-auc-graph.png")

    # ==========================================
    # SLIDE 24: Figure: Prediction Samples (DEDICATED SLIDE)
    # ==========================================
    add_full_image_slide(prs, "Figure 7: Qualitative Predictions on 20 Test Images", "docu/figures/20_sample_predictions.png")

    # ==========================================
    # SLIDE 25: Chapter 5: Conclusions
    # ==========================================
    slide_25 = create_base_slide(prs, "Chapter 5: Conclusions")
    
    tf_l25 = slide_25.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l25.word_wrap = True
    
    add_bullet(tf_l25, "AlexNet Performance: Achieved a classification accuracy of 98.91% on the test dataset.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l25, "Optimization efficiency: Pre-trained feature freezing minimized parameters, enabling rapid optimization.")
    add_bullet(tf_l25, "Zero Bias Imbalance Resilience: Stratification successfully prevented bias towards majority categories, yielding F1 score of 0.9910.")
    add_bullet(tf_l25, "Comparative Success: Yields competitive scores vs VGG/ResNets at a fraction of training times.")

    # ==========================================
    # SLIDE 26: Chapter 5: Recommendations & Future Work
    # ==========================================
    slide_26 = create_base_slide(prs, "Chapter 5: Recommendations & Future Work")
    
    tf_l26 = slide_26.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0)).text_frame
    tf_l26.word_wrap = True
    
    add_bullet(tf_l26, "Unfreezing schedules: Explore gradual layer unfreezing to fine-tune early kernels for sign boundaries.", size=18, bold=True, color=RGBColor(21, 57, 107))
    add_bullet(tf_l26, "Advanced synthetic augmentation: Use generative models to produce blurred and occluded training signs.")
    add_bullet(tf_l26, "Edge compiler conversion: Compile to TensorRT or ONNX for edge processors (Jetson/Pi).")
    add_bullet(tf_l26, "Dynamic ADAS Integration: Benchmark inference latency to ensure speeds match driver assistance systems.")
    
    # Save Presentation
    prs.save("docu/gtsrb_alexnet_presentation.pptx")
    print("New presentation saved successfully at docu/gtsrb_alexnet_presentation.pptx")

if __name__ == "__main__":
    main()
